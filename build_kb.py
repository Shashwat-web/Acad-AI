from pathlib import Path
from pypdf import PdfReader
from pptx import Presentation
from tqdm import tqdm
from sentence_transformers import SentenceTransformer
import numpy as np

DOCS_DIR = Path("knowledge_base/documents")

all_documents = []

def read_pdf(path):
    text = ""
    try:
        reader = PdfReader(path)
        for page in reader.pages:
            text += page.extract_text() or ""
    except Exception as e:
        print(f"PDF Error: {path} -> {e}")
    return text

def read_pptx(path):
    text = ""
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"PPTX Error: {path} -> {e}")
    return text

def read_txt(path):
    try:
        return open(path, "r", encoding="utf-8", errors="ignore").read()
    except Exception as e:
        print(f"TXT Error: {path} -> {e}")
        return ""
    
def chunk_text(text, chunk_size=1000, overlap=200):
    chunks = []

    start = 0

    while start < len(text):
        end = start + chunk_size

        chunks.append(text[start:end])

        start += chunk_size - overlap

    return chunks

for file in tqdm(list(DOCS_DIR.rglob("*"))):

    if file.suffix.lower() == ".pdf":
        text = read_pdf(file)

    elif file.suffix.lower() == ".pptx":
        text = read_pptx(file)

    elif file.suffix.lower() == ".txt":
        text = read_txt(file)

    else:
        continue

    if text.strip():

        chunks = chunk_text(text)

        for chunk in chunks:
            all_documents.append({
                "source": str(file),
                "text": chunk
            })

print(f"\nLoaded Documents: {len(all_documents)}")
print(f"Loaded Documents: {len(set(d['source'] for d in all_documents))}")
print(f"Total Chunks: {len(all_documents)}")

print(f"\nLoaded Documents: {len(set(d['source'] for d in all_documents))}")
print(f"Total Chunks: {len(all_documents)}")

print("\nLoading embedding model...")

model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")

# -------------------------
# Clean the chunks
# -------------------------

texts = []

for doc in all_documents:

    text = doc.get("text", "")

    if text is None:
        continue

    # Convert to string
    text = str(text)

    # Remove null characters
    text = text.replace("\x00", " ")

    # Remove extra spaces
    text = " ".join(text.split())

    if len(text) < 20:
        continue

    texts.append(text)

print(f"Chunks used for embeddings: {len(texts)}")

# Safety check
for i, t in enumerate(texts):
    if not isinstance(t, str):
        print(f"Invalid type at {i}: {type(t)}")
        raise TypeError("Non-string found!")

print("\nGenerating embeddings...")

embeddings = model.encode(
    texts,
    batch_size=16,
    show_progress_bar=True,
    convert_to_numpy=True,
    normalize_embeddings=True
)

print("\nEmbedding matrix shape:", embeddings.shape)