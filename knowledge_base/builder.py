"""
AcadAI - Academic Knowledge Base Builder

This script builds the permanent Academic Knowledge Base.

Workflow
--------
knowledge_base/documents
        ↓
Extract Text
        ↓
Create Chunk Objects
        ↓
Generate Embeddings
        ↓
Build FAISS Index
        ↓
knowledge_base/vector_store
    ├── index.faiss
    └── chunks.pkl

Run:
    uv run knowledge_base/builder.py
"""
import faiss
import pickle
import numpy as np

from retrieval.embeddings import get_embedding_model
from config import DEFAULT_EMBEDDING_MODEL
from pathlib import Path
from models.chunk import Chunk
# -------------------------------
# Configuration
# -------------------------------


DOCUMENTS_DIR = Path("knowledge_base/documents")
VECTOR_STORE_DIR = Path("knowledge_base/vector_store")

CHUNK_SIZE = 700
CHUNK_OVERLAP = 120

SUPPORTED_FILES = {
    ".pdf",
    ".ppt",
    ".pptx",
    ".txt",
}

from pypdf import PdfReader
from pptx import Presentation


def read_pdf(path: Path) -> str:
    text = ""

    try:
        reader = PdfReader(path)

        for page in reader.pages:
            text += page.extract_text() or ""

    except Exception as e:
        print(f"Skipped PDF: {path.name} -> {e}")

    return text


def read_pptx(path: Path) -> str:
    text = ""

    try:
        prs = Presentation(path)

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    except Exception as e:
        print(f"Skipped PPTX: {path.name} -> {e}")

    return text


def read_txt(path: Path) -> str:

    try:
        return path.read_text(
            encoding="utf-8",
            errors="ignore"
        )

    except Exception as e:
        print(f"Skipped TXT: {path.name} -> {e}")

        return ""

def chunk_text(
    text: str,
    source: str,
    chunk_size: int,
    overlap: int,
):
    """
    Split text into Chunk objects.
    """

    chunks = []

    start = 0
    chunk_id = 1

    text = " ".join(text.split())

    while start < len(text):

        end = start + chunk_size

        piece = text[start:end]

        if len(piece.strip()) > 50:

            chunks.append(
                Chunk(
                    doc_id=f"{Path(source).stem}::{chunk_id}",
                    source=Path(source).name,
                    page=0,
                    text=piece,
                )
            )

        start += chunk_size - overlap
        chunk_id += 1

    return chunks


def create_chunks(text: str, source: Path):
    """
    Convert a document into Chunk objects.
    """

    text = " ".join(text.split())

    chunks = []

    start = 0
    chunk_no = 1

    while start < len(text):

        piece = text[start:start + CHUNK_SIZE]

        piece = piece.encode(
            "utf-16",
            "surrogatepass"
        ).decode(
            "utf-16",
            "ignore"
        )

        piece = piece.encode(
            "utf-8",
            "ignore"
        ).decode(
            "utf-8",
            "ignore"
        )

        if len(piece.strip()) >= 50:

            chunks.append(
                Chunk(
                    doc_id=f"{source.stem}::{chunk_no}",
                    source=source.name,
                    page=0,
                    text=piece,
                )
            )

        start += CHUNK_SIZE - CHUNK_OVERLAP
        chunk_no += 1

    return chunks

def main():

    print("=" * 60)
    print("AcadAI Knowledge Base Builder")
    print("=" * 60)

    # Step 1
    print("\n[1/5] Scanning documents...")

    files = []

    for ext in SUPPORTED_FILES:
        files.extend(DOCUMENTS_DIR.rglob(f"*{ext}"))

    files = sorted(files)

    print(f"Found {len(files)} documents")

    if not files:
        print("No documents found.")
        return
    


    print("[2/5] Extracting text...")

    documents = []

    for file in files:

        suffix = file.suffix.lower()

        if suffix == ".pdf":
            text = read_pdf(file)

        elif suffix == ".pptx":
            text = read_pptx(file)

        elif suffix == ".ppt":
            print(f"Skipped PPT: {file.name} (Legacy .ppt format)")
            continue

        elif suffix == ".txt":
            text = read_txt(file)

        else:
            continue

        if text.strip():

            documents.append(
                {
                    "path": file,
                    "text": text,
                }
            )

    print(f"Documents extracted: {len(documents)}")

    # Step 3
    print("[3/5] Creating chunks...")

    all_chunks = []

    for doc in documents:

        chunks = create_chunks(
            doc["text"],
            doc["path"],
        )

        all_chunks.extend(chunks)

    print(f"Chunks created: {len(all_chunks)}")

    # Step 4
    # Step 4
    print("[4/5] Generating embeddings...")

    model = get_embedding_model(DEFAULT_EMBEDDING_MODEL)

    if model is None:
        raise RuntimeError("Failed to load embedding model.")

    texts = []

    for i, chunk in enumerate(all_chunks):

        if not isinstance(chunk.text, str):
            print(f"\nInvalid chunk at index {i}")
            print("Type:", type(chunk.text))
            print(chunk)
            break

        texts.append(chunk.text)

    embeddings = []

    for i, text in enumerate(texts):

        try:

            emb = model.encode(
                text,
                convert_to_numpy=True,
                normalize_embeddings=True,
            )

            embeddings.append(emb)

        except Exception as e:

            print("\nFAILED CHUNK:", i)
            print("Python type:", type(text))
            print("Length:", len(text))
            print("Representation:")
            print(repr(text))
            print("\nException:", e)

            raise

    embeddings = np.asarray(
        embeddings,
        dtype=np.float32,
    )    

    print(f"Embeddings shape: {embeddings.shape}")

    # Step 5
    print("[5/5] Building FAISS index...")

    VECTOR_STORE_DIR.mkdir(
        parents=True,
        exist_ok=True,
    )

    index = faiss.IndexFlatIP(
        embeddings.shape[1]
    )

    index.add(embeddings)

    faiss.write_index(
        index,
        str(VECTOR_STORE_DIR / "index.faiss"),
    )

    with open(
        VECTOR_STORE_DIR / "chunks.pkl",
        "wb",
    ) as f:
        pickle.dump(all_chunks, f)

    print(f"FAISS vectors : {index.ntotal}")
    print(f"Chunks saved  : {len(all_chunks)}")
    print(f"Location      : {VECTOR_STORE_DIR}")

print("\nKnowledge Base Build Complete")


if __name__ == "__main__":
    main()